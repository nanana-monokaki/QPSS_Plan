import os
import requests
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def create_collage():
    # 1. Download images
    images = []
    with open('urls.txt', 'r') as f:
        urls = [line.strip() for line in f if line.strip()]
    
    for url in urls[:9]:  # Download first 9 images
        print(f"Downloading {url}...")
        try:
            response = requests.get(url)
            if response.status_code == 200:
                img = Image.open(BytesIO(response.content)).convert("RGBA")
                images.append(img)
        except Exception as e:
            print(f"Failed to download {url}: {e}")
    
    if not images:
        print("No images found.")
        return None
        
    # 2. Create a 3x3 grid canvas for the collage
    # Standard 16:9 ratio, high res (1920x1080)
    canvas_w, canvas_h = 1920, 1080
    collage = Image.new('RGBA', (canvas_w, canvas_h), (255, 255, 255, 0))
    
    cell_w = canvas_w // 3
    cell_h = canvas_h // 3
    
    idx = 0
    for row in range(3):
        for col in range(3):
            if idx < len(images):
                img = images[idx]
                
                # Resize and crop to fill the cell (Aspect Fill)
                img_ratio = img.width / img.height
                cell_ratio = cell_w / cell_h
                
                if img_ratio > cell_ratio:
                    # Image is wider
                    new_h = cell_h
                    new_w = int(img_ratio * cell_h)
                    img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                    # Crop center
                    left = (new_w - cell_w) / 2
                    img = img.crop((left, 0, left + cell_w, cell_h))
                else:
                    # Image is taller
                    new_w = cell_w
                    new_h = int(cell_w / img_ratio)
                    img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                    # Crop top
                    top = (new_h - cell_h) / 2
                    img = img.crop((0, top, cell_w, top + cell_h))
                    
                # Paste into the collage
                x = col * cell_w
                y = row * cell_h
                collage.paste(img, (x, y))
                idx += 1
                
    # 3. Apply 50% transparency
    print("Applying 50% transparency...")
    r, g, b, a = collage.split()
    # 127 is approx 50% of 255
    alpha = a.point(lambda p: 127 if p > 0 else 0)
    collage.putalpha(alpha)
    
    collage_path = "collage_bg.png"
    collage.save(collage_path)
    print(f"Collage saved to {collage_path}")
    return collage_path

def update_pptx(collage_path):
    pptx_path = "聞き耳アワーシリーズ企画書_修正版.pptx"
    print(f"Opening presentation: {pptx_path}")
    prs = Presentation(pptx_path)
    
    # Target 7th slide (index 6)
    slide_index = 6
    if len(prs.slides) <= slide_index:
        print(f"Error: Slide {slide_index + 1} does not exist.")
        return
        
    slide = prs.slides[slide_index]
    
    # Add picture to slide
    print("Adding picture to slide 7...")
    pic = slide.shapes.add_picture(collage_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # Send picture to back
    # spTree is the XML element containing all shapes. The first shape is usually at index 2 (after p:bg and p:cSld).
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    
    output_path = "聞き耳アワーシリーズ企画書_修正版.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    collage_img = create_collage()
    if collage_img:
        update_pptx(collage_img)
