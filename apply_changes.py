import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

file_path = os.path.join('d:\\\\', 'OneDrive', 'クリエイティブ関連', 'QPSS', 'QPSS_KikimimiSeries', 'Plan', '聞き耳アワーシリーズ企画書_鴨FB対応_260228.pptx')
output_path = os.path.join('d:\\\\', 'OneDrive', 'クリエイティブ関連', 'QPSS', 'QPSS_KikimimiSeries', 'Plan', '聞き耳アワーシリーズ企画書_調整版_03.pptx')

prs = Presentation(file_path)

# -----------------
# Slide 3 (Index 2)
# -----------------
slide3 = prs.slides[2]
# Shape 3: 直木賞作家・姫野カオルコがプロデュース！ -> 直木賞作家・姫野カオルコが「耳で読む文学」を厳選チョイス＆完全プロデュース！
shape3_3 = slide3.shapes[3]
shape3_3.text = "直木賞作家・姫野カオルコが「耳で読む文学」を厳選チョイス＆完全プロデュース！"

# Shape 3: 新しい文学体験型...
shape3_4 = slide3.shapes[4]
shape3_4.text = (
    "選び抜かれた珠玉の物語が、いま新しい「文学体験型朗読劇」として生まれ変わります。\n"
    "劇場という空間で声・音・映像とともにQPSSが立体化し、\n"
    "さらにライブ体験を起点としてポッドキャスト・音声配信へと展開。\n"
    "「耳から出会う文学」をあなたの日常へ届けます。"
)

# -----------------
# Slide 9 (Index 8)
# -----------------
slide9 = prs.slides[8]
shape9_4 = slide9.shapes[4]
# Adjust font size to fit new text length
shape9_4.text = (
    "映画・ドラマの企画開発から、ゲーム、企業キャラクター開発まで、構成力と創造性であらゆる物語を設計・構築するクリエイティブ集団です。\n\n"
    "【 圧倒的没入感！脚本家チームがつくる「耳で聴く演劇舞台」 】\n"
    "本シリーズでは、QPSSのプロ脚本家チームが脚本構成・舞台演出設計・音響設計すべてを統括。\n"
    "単なる朗読ではなく、文学作品を“舞台言語”へ見事に翻訳し、上質な「耳で聴く演劇舞台」空間を創り上げます。\n\n"
    "HP：http://www.quobo-pic.com"
)
for paragraph in shape9_4.text_frame.paragraphs:
    paragraph.font.size = Pt(14)
    if "【" in paragraph.text:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(192, 0, 0) # Dark red for emphasis

# -----------------
# Slide 11 (Index 10)
# -----------------
slide11 = prs.slides[10]
shape11_5 = slide11.shapes[5]
# Original text: 
# 【第１弾】 「エンドレス・ラブ」（徳間書店）11月または12月想定
# 【第２弾】 「X博士」
# 【第３弾】 「探偵物語」

# Move shape 5 down and narrow it to fit only the future lineup
shape11_5.top = int(shape11_5.top + Inches(1.0))
shape11_5.text = (
    "今後のラインナップ予定：\n"
    "・【第２弾】 「X博士」\n"
    "・【第３弾】 「探偵物語」"
)
for paragraph in shape11_5.text_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(128, 128, 128) # Gray text

# Now add a new prominent shape for the 1st installment
from pptx.enum.shapes import MSO_SHAPE
left = int(slide11.shapes[4].left) # Same left as "今後の公演予定" title
top = int(slide11.shapes[4].top) + int(Inches(0.6))
width = int(Inches(7.0))
height = int(Inches(1.0))

new_shape = slide11.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)
# Style the shape
new_shape.fill.solid()
new_shape.fill.fore_color.rgb = RGBColor(255, 242, 204) # Light yellow background
new_shape.line.color.rgb = RGBColor(255, 192, 0) # Orange/Gold outline
new_shape.line.width = Pt(2)

# Set text
tf = new_shape.text_frame
tf.text = "★ 第1弾 注目公演"
p0 = tf.paragraphs[0]
p0.font.bold = True
p0.font.size = Pt(14)
p0.font.color.rgb = RGBColor(192, 0, 0)

p1 = tf.add_paragraph()
p1.text = "「エンドレス・ラブ」（徳間書店） 11月または12月想定"
p1.font.bold = True
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(0, 0, 0)

prs.save(output_path)
print(f"Saved {output_path}")
