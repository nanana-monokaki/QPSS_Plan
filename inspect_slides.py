import sys
import os
from pptx import Presentation

file_path = os.path.join('d:\\\\', 'OneDrive', 'クリエイティブ関連', 'QPSS', 'QPSS_KikimimiSeries', 'Plan', '聞き耳アワーシリーズ企画書_鴨FB対応_260228.pptx')
prs = Presentation(file_path)

with open('slide_structure.txt', 'w', encoding='utf-8') as f:
    for slide_num in [2, 8, 10]:  # 0-indexed: P3, P9, P11
        slide = prs.slides[slide_num]
        f.write(f'--- Slide {slide_num + 1} ---\n')
        for i, shape in enumerate(slide.shapes):
            f.write(f'Shape {i}: type={shape.shape_type} name="{shape.name}" left={shape.left} top={shape.top} width={shape.width} height={shape.height}\n')
            try:
                if shape.has_text_frame:
                    text = shape.text.replace('\n', '\\n')
                    f.write(f'  Text: {text[:200]}...\n')
                elif shape.has_table:
                    f.write('  Table\n')
            except Exception as e:
                f.write(f'  Error: {e}\n')
