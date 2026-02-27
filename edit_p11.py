from pptx import Presentation

def edit_p11():
    pptx_path = "聞き耳アワーシリーズ企画書_調整版_02.pptx"
    out_path = "聞き耳アワーシリーズ企画書_調整版_03.pptx"
    prs = Presentation(pptx_path)
    
    # 対象スライド：P11 (インデックス10)
    slide = prs.slides[10]
    
    perf_shape = None
    cast_shape = None
    
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
            if "公演回数" in shape.text:
                perf_shape = shape
                print(f"Found 公演回数 at idx returning top {shape.top}")
            elif "キャスト数" in shape.text:
                cast_shape = shape
                print(f"Found キャスト数 at idx returning top {shape.top}")
                
    if perf_shape and cast_shape:
        # 文字列置換
        for paragraph in perf_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "公演回数" in run.text:
                    run.text = run.text.replace("公演回数", "ステージ数")
                    print("Replaced text '公演回数' -> 'ステージ数'")
                    
        # Y座標(top)入れ替え
        temp_top = perf_shape.top
        perf_shape.top = cast_shape.top
        cast_shape.top = temp_top
        
        print("Swapped positions.")
        prs.save(out_path)
        print(f"Saved to {out_path}")
    else:
        print("Shapes not found.")

if __name__ == "__main__":
    edit_p11()
