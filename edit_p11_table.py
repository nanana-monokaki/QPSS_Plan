import copy
from pptx import Presentation

def reorder_table_rows():
    pptx_path = "聞き耳アワーシリーズ企画書_調整版_02.pptx"
    out_path = "聞き耳アワーシリーズ企画書_調整版_03.pptx"
    prs = Presentation(pptx_path)
    slide = prs.slides[10]
    
    target_table = None
    for shape in slide.shapes:
        if shape.has_table:
            target_table = shape.table
            break
            
    if not target_table:
        print("Table not found")
        return
        
    row2 = target_table.rows[2] # キャスト数
    row3 = target_table.rows[3] # 客席数
    row4 = target_table.rows[4] # 公演回数
    
    # Store text
    text_r2_c0 = row2.cells[0].text
    text_r2_c1 = row2.cells[1].text
    text_r3_c0 = row3.cells[0].text
    text_r3_c1 = row3.cells[1].text
    text_r4_c0 = row4.cells[0].text.replace("公演回数", "ステージ数")
    text_r4_c1 = row4.cells[1].text
    
    # helper for fast replacement without losing format completely (we overwrite the first run)
    def set_cell_text(cell, text):
        if cell.text_frame.paragraphs:
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                # keep first run formatting
                p.runs[0].text = text
                # remove other runs
                for i in range(len(p.runs)-1, 0, -1):
                    # removing runs from python-pptx is not trivial via API,
                    # but we can try clearing the text.
                    p.runs[i].text = ""
            else:
                cell.text = text
        else:
            cell.text = text

    # New Row 2 is Old Row 4
    set_cell_text(row2.cells[0], text_r4_c0)
    set_cell_text(row2.cells[1], text_r4_c1)
    
    # New Row 3 is Old Row 2
    set_cell_text(row3.cells[0], text_r2_c0)
    set_cell_text(row3.cells[1], text_r2_c1)

    # New Row 4 is Old Row 3
    set_cell_text(row4.cells[0], text_r3_c0)
    set_cell_text(row4.cells[1], text_r3_c1)
    
    prs.save(out_path)
    print(f"Saved to {out_path}")

if __name__ == "__main__":
    reorder_table_rows()
