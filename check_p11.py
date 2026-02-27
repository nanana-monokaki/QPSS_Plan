from pptx import Presentation
prs = Presentation('聞き耳アワーシリーズ企画書_調整版_03.pptx')
slide = prs.slides[10]
for shape in slide.shapes:
    if shape.has_table:
        for r_idx, row in enumerate(shape.table.rows):
            texts = [cell.text.replace('\n', '') for cell in row.cells]
            print(f'{r_idx}: {texts}')
