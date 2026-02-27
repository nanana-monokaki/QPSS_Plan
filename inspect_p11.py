from pptx import Presentation

def inspect_slide():
    pptx_path = "聞き耳アワーシリーズ企画書_調整版_02.pptx"
    prs = Presentation(pptx_path)
    
    # 対象スライド：P11 (インデックス10)
    slide_index = 10
    slide = prs.slides[slide_index]
    
    print("--- Shapes on Slide 11 ---")
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
            text = shape.text.replace("\n", " ")
            print(f"[{i}] top: {shape.top}, left: {shape.left} | text: {text[:50]}")
            
if __name__ == "__main__":
    inspect_slide()
