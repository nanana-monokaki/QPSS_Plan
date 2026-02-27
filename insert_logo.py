from pptx import Presentation
from pptx.util import Inches

def insert_logo():
    pptx_path = "聞き耳アワーシリーズ企画書_調整版_03.pptx"
    out_path = "聞き耳アワーシリーズ企画書_調整版_04.pptx"
    logo_path = "QPSSTOP_transparent.png"
    
    prs = Presentation(pptx_path)
    slide = prs.slides[11] # P12
    
    # 画像のサイズを指定（幅2インチ）
    logo_width = Inches(2.0)
    
    # スライドの右下あたりに配置する
    # 右下余白を0.5インチと想定
    margin = Inches(0.5)
    left = prs.slide_width - logo_width - margin
    # 高さはアスペクト比で自動計算されるので一旦適当なyに置く
    # 後から画像の元の比率に合わせて高さを決め、bottom合わせにするのが理想だが、
    # widthだけ指定すると自動でheightも決まるので、topの位置を適当にして後で微調整するか、
    # 概ね下部に配置されるようにYを固定する。
    top = prs.slide_height - Inches(1.5) - margin
    
    slide.shapes.add_picture(logo_path, left, top, width=logo_width)
    
    prs.save(out_path)
    print(f"Saved to {out_path}")

if __name__ == "__main__":
    insert_logo()
