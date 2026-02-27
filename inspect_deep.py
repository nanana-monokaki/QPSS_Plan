from pptx import Presentation

def inspect_table():
    pptx_path = "聞き耳アワーシリーズ企画書_調整版_02.pptx"
    prs = Presentation(pptx_path)
    slide = prs.slides[10]
    
    for shape in slide.shapes:
        if shape.has_table:
            print(f"Table id: {shape.shape_id}")
            for r_idx, row in enumerate(shape.table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.replace('\n', ' '))
                print(f"  Row {r_idx}: {row_text}")

if __name__ == "__main__":
    inspect_table()
